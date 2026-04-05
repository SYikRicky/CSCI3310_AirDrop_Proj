#!/usr/bin/env python3
"""Generate the CSCI3310 DroidDrop demo presentation — light theme with syntax highlighting."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import re

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Light colour palette ────────────────────────────────────────────────────
BG_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BG_LIGHT  = RGBColor(0xF5, 0xF7, 0xFA)
BG_CARD   = RGBColor(0xFF, 0xFF, 0xFF)
BORDER    = RGBColor(0xE0, 0xE4, 0xEB)

BLUE      = RGBColor(0x1A, 0x73, 0xE8)   # primary
INDIGO    = RGBColor(0x5C, 0x6B, 0xC0)
TEAL      = RGBColor(0x00, 0x89, 0x7B)
GREEN     = RGBColor(0x2E, 0x7D, 0x32)
ORANGE    = RGBColor(0xE6, 0x5C, 0x00)
RED       = RGBColor(0xC6, 0x28, 0x28)
PINK      = RGBColor(0xAD, 0x14, 0x57)

BLACK     = RGBColor(0x20, 0x24, 0x2A)
DARK      = RGBColor(0x3C, 0x40, 0x43)
GRAY      = RGBColor(0x5F, 0x63, 0x68)
LIGHT_TXT = RGBColor(0x80, 0x86, 0x8B)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)

# Syntax highlight colours (VS Code light-like)
SYN_KEYWORD = RGBColor(0xAF, 0x00, 0xDB)  # purple
SYN_STRING  = RGBColor(0xA3, 0x11, 0x15)  # dark red
SYN_TYPE    = RGBColor(0x26, 0x7F, 0x99)  # teal
SYN_COMMENT = RGBColor(0x6A, 0x99, 0x55)  # green
SYN_METHOD  = RGBColor(0x79, 0x56, 0x48)  # brown
SYN_PLAIN   = RGBColor(0x33, 0x33, 0x33)
SYN_NUMBER  = RGBColor(0x09, 0x86, 0x58)  # green-ish
SYN_SYMBOL  = RGBColor(0x00, 0x00, 0x00)

CODE_BG     = RGBColor(0xF8, 0xF8, 0xF8)

JAVA_KEYWORDS = {
    'extends', 'implements', 'interface', 'class', 'public', 'private',
    'protected', 'static', 'final', 'void', 'return', 'new', 'import',
    'package', 'if', 'else', 'for', 'while', 'try', 'catch', 'throw',
    'throws', 'abstract', 'boolean', 'int', 'long', 'double', 'float',
    'String', 'true', 'false', 'null', 'enum', 'this', 'super',
}

JAVA_TYPES = {
    'AppCompatActivity', 'Fragment', 'LifecycleService', 'RecyclerView',
    'Adapter', 'ViewHolder', 'Parcelable', 'Intent', 'Uri', 'Bundle',
    'List', 'Map', 'Set', 'HashMap', 'ArrayList', 'SharedPreferences',
    'JSONArray', 'JSONObject', 'MediaPlayer', 'MediaRecorder',
    'ConnectionsClient', 'Payload', 'PayloadCallback', 'Context',
    'File', 'FileMetadata', 'DeviceInfo', 'TransferProgress',
    'ChatMessage', 'SharedFile', 'Notification', 'NotificationChannel',
    'FirebaseFirestore', 'FirebaseStorage', 'MapView', 'GeoPoint',
    'ContentResolver', 'MediaStore', 'FusedLocationProviderClient',
    'ExecutorService', 'Handler', 'Looper', 'View', 'ImageView',
    'Dialog', 'AlertDialog', 'FloatingActionButton', 'BottomNavigationView',
    'MaterialCardView', 'ProgressBar', 'TextView', 'ChatAdapter',
    'DeviceAdapter', 'FileListAdapter', 'HistoryPeerAdapter',
    'NearbyConnectionsManager', 'ConnectionListener', 'TransferListener',
    'ChatListener', 'FileTransferService', 'FileStorageManager',
    'ChatHistoryManager', 'SharedDriveRepository', 'TransferEventBus',
    'ChatFileReceivedCallback', 'OnFileSavedCallback', 'POJO',
}


def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_card(slide, left, top, width, height, fill=BG_CARD, border_color=BORDER):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1)
    shape.shadow.inherit = False
    if hasattr(shape, 'adjustments') and len(shape.adjustments) > 0:
        shape.adjustments[0] = 0.03
    return shape

def add_text(slide, left, top, width, height, text, font_size=18, color=BLACK,
             bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=15,
                    color=DARK, bullet_color=BLUE):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(3)
        p.space_after = Pt(2)
        b = p.add_run()
        b.text = "  \u25B8  "
        b.font.size = Pt(font_size)
        b.font.color.rgb = bullet_color
        b.font.bold = True
        b.font.name = "Calibri"
        r = p.add_run()
        r.text = item
        r.font.size = Pt(font_size)
        r.font.color.rgb = color
        r.font.name = "Calibri"
    return txBox

def _tokenize_java(line):
    """Split a Java-like line into (text, token_type) pairs for syntax highlighting."""
    tokens = []
    # Handle comment lines
    stripped = line.lstrip()
    if stripped.startswith('//') or stripped.startswith('*') or stripped.startswith('/*'):
        return [(line, 'comment')]

    # Regex-based simple tokenizer
    pattern = re.compile(
        r'("(?:[^"\\]|\\.)*")'       # strings
        r'|(\b\d+(?:\.\d+)?\b)'      # numbers
        r"|(\b[A-Za-z_]\w*\b)"       # identifiers
        r'|(\s+)'                     # whitespace
        r'|(.)'                       # other chars
    )
    for m in pattern.finditer(line):
        s, num, ident, ws, other = m.groups()
        if s:
            tokens.append((s, 'string'))
        elif num:
            tokens.append((num, 'number'))
        elif ident:
            if ident in JAVA_KEYWORDS:
                tokens.append((ident, 'keyword'))
            elif ident in JAVA_TYPES or (ident[0].isupper() and len(ident) > 1):
                tokens.append((ident, 'type'))
            else:
                tokens.append((ident, 'plain'))
        elif ws:
            tokens.append((ws, 'plain'))
        elif other:
            tokens.append((other, 'symbol'))
    return tokens

TOKEN_COLORS = {
    'keyword': SYN_KEYWORD,
    'string':  SYN_STRING,
    'type':    SYN_TYPE,
    'comment': SYN_COMMENT,
    'number':  SYN_NUMBER,
    'symbol':  SYN_SYMBOL,
    'plain':   SYN_PLAIN,
}

def add_code_block(slide, left, top, width, height, code, font_size=12):
    """Add a code block with syntax highlighting."""
    add_card(slide, left, top, width, height, fill=CODE_BG, border_color=BORDER)
    txBox = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.12),
                                      width - Inches(0.4), height - Inches(0.24))
    tf = txBox.text_frame
    tf.word_wrap = True
    lines = code.split('\n')
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(1)
        p.space_after = Pt(1)
        tokens = _tokenize_java(line)
        if not tokens:
            r = p.add_run()
            r.text = " "
            r.font.size = Pt(font_size)
            r.font.name = "Consolas"
            continue
        for text, tok_type in tokens:
            r = p.add_run()
            r.text = text
            r.font.size = Pt(font_size)
            r.font.name = "Consolas"
            r.font.color.rgb = TOKEN_COLORS.get(tok_type, SYN_PLAIN)
            r.font.bold = tok_type == 'keyword'
    return txBox

# Accent stripe at top of a slide
def add_top_stripe(slide, color=BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.08))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_WHITE)
# large accent rectangle at top
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(3.2))
shape.fill.solid()
shape.fill.fore_color.rgb = BLUE
shape.line.fill.background()

add_text(slide, Inches(1), Inches(0.8), Inches(11), Inches(1),
         "DroidDrop", font_size=64, color=WHITE, bold=True)
add_text(slide, Inches(1), Inches(2.0), Inches(11), Inches(0.7),
         "Peer-to-Peer File Sharing & Chat for Android", font_size=26, color=RGBColor(0xBB, 0xDE, 0xFB))

add_text(slide, Inches(1), Inches(3.8), Inches(11), Inches(0.6),
         "CSCI3310 Mobile Computing  |  Spring 2026", font_size=22, color=DARK)

# Tech pills
techs = ["Nearby Connections", "Firebase", "OSMDroid", "Material Design 3", "Java"]
x = Inches(1)
for tech in techs:
    w = Inches(len(tech) * 0.12 + 0.4)
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(4.8), w, Inches(0.45))
    pill.fill.solid()
    pill.fill.fore_color.rgb = BG_LIGHT
    pill.line.color.rgb = BORDER
    pill.line.width = Pt(1)
    if hasattr(pill, 'adjustments') and len(pill.adjustments) > 0:
        pill.adjustments[0] = 0.5
    add_text(slide, x + Inches(0.15), Inches(4.82), w - Inches(0.3), Inches(0.4),
             tech, font_size=14, color=BLUE, bold=True, alignment=PP_ALIGN.CENTER)
    x += w + Inches(0.15)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Architecture Overview
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_LIGHT)
add_top_stripe(slide)

add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         "Architecture Overview", font_size=36, color=BLACK, bold=True)
add_text(slide, Inches(0.6), Inches(0.9), Inches(12), Inches(0.4),
         "MVC Pattern  |  Callback Interfaces  |  Event-Driven Communication", font_size=17, color=GRAY)

pillars = [
    ("P2P Chat & Transfer", [
        "Google Nearby Connections API",
        "P2P_CLUSTER strategy",
        "BYTES + FILE payloads",
        "Device discovery & advertising",
    ], BLUE),
    ("Cloud Shared Drive", [
        "Firebase Cloud Storage",
        "Cloud Firestore (real-time)",
        "Snapshot listeners",
        "Upload / Download / Delete",
    ], INDIGO),
    ("Offline Map & Location", [
        "OSMDroid (OpenStreetMap)",
        "OSRM walking routes",
        "GPS location sharing",
        "FusedLocationProviderClient",
    ], TEAL),
]
for i, (title, items, color) in enumerate(pillars):
    left = Inches(0.5 + i * 4.2)
    add_card(slide, left, Inches(1.6), Inches(3.9), Inches(2.9))
    # color bar at top of card
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(1.6), Inches(3.9), Inches(0.06))
    bar.fill.solid(); bar.fill.fore_color.rgb = color; bar.line.fill.background()
    add_text(slide, left + Inches(0.2), Inches(1.75), Inches(3.5), Inches(0.4),
             title, font_size=19, color=color, bold=True)
    add_bullet_list(slide, left + Inches(0.15), Inches(2.15), Inches(3.5), Inches(2.0),
                    items, font_size=14, color=DARK, bullet_color=color)

# Tech stack bar
add_card(slide, Inches(0.5), Inches(4.85), Inches(12.3), Inches(0.65))
add_text(slide, Inches(0.7), Inches(4.9), Inches(12), Inches(0.55),
         "Java  |  Min SDK 24  |  Target SDK 36  |  play-services-nearby 18.7.0  |  "
         "Firebase BOM  |  OSMDroid  |  Material 3  |  Lifecycle-Aware Components",
         font_size=13, color=GRAY)

# Class hierarchy summary
add_card(slide, Inches(0.5), Inches(5.7), Inches(12.3), Inches(1.5))
add_text(slide, Inches(0.7), Inches(5.75), Inches(4), Inches(0.35),
         "Class Hierarchy Summary", font_size=16, color=BLACK, bold=True)
add_code_block(slide, Inches(0.7), Inches(6.1), Inches(11.8), Inches(1.0),
    "MainActivity (AppCompatActivity) --> 5 Fragments + 4 Adapters\n"
    "NearbyConnectionsManager --> TransferEventBus (3 callback interfaces)\n"
    "FileTransferService (LifecycleService) | SharedDriveRepository (Firebase)\n"
    "5 Models: DeviceInfo, FileMetadata (Parcelable), TransferProgress (Parcelable), ChatMessage, SharedFile",
    font_size=12)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Activities
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_LIGHT)
add_top_stripe(slide)

add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         "Activities (2)", font_size=36, color=BLACK, bold=True)

# MainActivity
add_card(slide, Inches(0.5), Inches(1.2), Inches(6.0), Inches(5.9))
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.2), Inches(6.0), Inches(0.06))
bar.fill.solid(); bar.fill.fore_color.rgb = BLUE; bar.line.fill.background()
add_text(slide, Inches(0.7), Inches(1.35), Inches(5.6), Inches(0.4),
         "MainActivity", font_size=22, color=BLUE, bold=True)
add_code_block(slide, Inches(0.7), Inches(1.8), Inches(5.6), Inches(1.4),
    "public class MainActivity extends AppCompatActivity\n"
    "    implements ConnectionListener,\n"
    "               TransferListener,\n"
    "               ChatListener,\n"
    "               OnFileSavedCallback,\n"
    "               ChatFileReceivedCallback",
    font_size=12)

add_text(slide, Inches(0.7), Inches(3.3), Inches(5.6), Inches(0.35),
         "Key Techniques:", font_size=15, color=BLACK, bold=True)
add_bullet_list(slide, Inches(0.7), Inches(3.6), Inches(5.4), Inches(3.3), [
    "FragmentContainerView + BottomNavigationView",
    "ActivityResultLauncher (permissions & file picker)",
    "FusedLocationProviderClient (GPS location)",
    "SharedPreferences (device name persistence)",
    "AlertDialog (chat invitations)",
    "configChanges (survive rotation)",
    "Intent.ACTION_SEND (cross-app share)",
], font_size=13, color=DARK, bullet_color=BLUE)

# MapActivity
add_card(slide, Inches(6.8), Inches(1.2), Inches(6.0), Inches(5.9))
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(1.2), Inches(6.0), Inches(0.06))
bar.fill.solid(); bar.fill.fore_color.rgb = TEAL; bar.line.fill.background()
add_text(slide, Inches(7.0), Inches(1.35), Inches(5.6), Inches(0.4),
         "MapActivity", font_size=22, color=TEAL, bold=True)
add_code_block(slide, Inches(7.0), Inches(1.8), Inches(5.6), Inches(0.5),
    "public class MapActivity extends AppCompatActivity",
    font_size=12)

add_text(slide, Inches(7.0), Inches(2.45), Inches(5.6), Inches(0.35),
         "Key Techniques:", font_size=15, color=BLACK, bold=True)
add_bullet_list(slide, Inches(7.0), Inches(2.75), Inches(5.4), Inches(4.2), [
    "OSMDroid MapView (OpenStreetMap tiles)",
    "MyLocationNewOverlay (current GPS position)",
    "GpsMyLocationProvider (location updates)",
    "Marker overlay (shared location pins)",
    "Polyline overlay (walking route)",
    "OSRM public API (route calculation)",
    "HttpURLConnection (network calls)",
    "ExecutorService (background thread)",
    "JSONObject / JSONArray (route parsing)",
    "Intent extras (EXTRA_LAT, EXTRA_LNG)",
], font_size=13, color=DARK, bullet_color=TEAL)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Fragments
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_LIGHT)
add_top_stripe(slide)

add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         "Fragments (5)", font_size=36, color=BLACK, bold=True)
add_text(slide, Inches(0.6), Inches(0.85), Inches(12), Inches(0.4),
         "All extend Fragment  |  Managed via FragmentManager + back stack  |  Arguments via Bundle",
         font_size=16, color=GRAY)

fragments = [
    ("FileListFragment", "Shared Drive", INDIGO, [
        "RecyclerView + FileListAdapter",
        "Firestore ListenerRegistration (real-time)",
        "FloatingActionButton (file picker)",
        "onDestroyView removes listener",
    ]),
    ("ChatRoomFragment", "Chat UI", BLUE, [
        "RecyclerView + ChatAdapter (2 view types)",
        "MediaRecorder (voice, audio/mp4)",
        "MediaPlayer (playback, prepareAsync)",
        "Inline ImageView (thumbnails)",
        "Location sharing + file attach buttons",
        "Connection status banner",
    ]),
    ("ChatDeviceListFragment", "Chat Lobby", TEAL, [
        "Dual RecyclerView layout:",
        "  DeviceAdapter (live devices)",
        "  HistoryPeerAdapter (past chats)",
        "ChatHistoryManager integration",
        "'Connect' on online history peers",
    ]),
    ("DeviceDiscoveryFragment", "File Transfer", ORANGE, [
        "RecyclerView + DeviceAdapter",
        "ProgressBar (transfer progress)",
        "Device click triggers file send",
    ]),
    ("SendModeFragment", "Firebase Upload", PINK, [
        "Firebase upload with progress bar",
        "ActivityResultLauncher (file picker)",
        "File metadata display",
    ]),
]

positions = [
    (Inches(0.35), Inches(1.45), Inches(4.0), Inches(2.5)),
    (Inches(4.55), Inches(1.45), Inches(4.3), Inches(2.5)),
    (Inches(9.05), Inches(1.45), Inches(4.0), Inches(2.5)),
    (Inches(0.35), Inches(4.25), Inches(6.15), Inches(2.1)),
    (Inches(6.7), Inches(4.25), Inches(6.35), Inches(2.1)),
]

for i, (name, subtitle, color, items) in enumerate(fragments):
    l, t, w, h = positions[i]
    add_card(slide, l, t, w, h)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, Inches(0.05))
    bar.fill.solid(); bar.fill.fore_color.rgb = color; bar.line.fill.background()
    add_text(slide, l + Inches(0.15), t + Inches(0.1), w - Inches(0.3), Inches(0.35),
             name, font_size=16, color=color, bold=True)
    add_text(slide, l + Inches(0.15), t + Inches(0.38), w - Inches(0.3), Inches(0.25),
             subtitle, font_size=12, color=GRAY)
    for j, item in enumerate(items):
        add_text(slide, l + Inches(0.2), t + Inches(0.65 + j * 0.27), w - Inches(0.35), Inches(0.27),
                 "\u25B8  " + item, font_size=11, color=DARK)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Nearby Connections API
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_LIGHT)
add_top_stripe(slide)

add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         "Nearby Connections API", font_size=36, color=BLACK, bold=True)
add_text(slide, Inches(0.6), Inches(0.85), Inches(12), Inches(0.4),
         "NearbyConnectionsManager.java  |  Strategy: P2P_CLUSTER  |  Service ID: com.example.dropdroid",
         font_size=15, color=GRAY)

# Left: Core API
add_card(slide, Inches(0.4), Inches(1.45), Inches(6.2), Inches(5.6))
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(1.45), Inches(6.2), Inches(0.05))
bar.fill.solid(); bar.fill.fore_color.rgb = BLUE; bar.line.fill.background()
add_text(slide, Inches(0.6), Inches(1.55), Inches(5.8), Inches(0.4),
         "Core API Calls", font_size=20, color=BLUE, bold=True)

methods = [
    ("Discovery", "startAdvertising() / startDiscovery()\nConnectionsClient from Nearby.getConnectionsClient()"),
    ("Connection", "requestConnection() / acceptConnection()\nConnectionLifecycleCallback (initiate, result, disconnect)"),
    ("Chat Text", 'Payload.fromBytes("CHAT|sender|ts|text")'),
    ("Location", 'Payload.fromBytes("LOCATION|sender|ts|lat|lng")'),
    ("Chat File", 'Payload.fromBytes("CHATFILE|sender|name|mime|size|"\n    + rawBytes)  // BYTES workaround for API 36'),
    ("Callbacks", "EndpointDiscoveryCallback (found / lost)\nPayloadCallback (received / transferUpdate)"),
]
y = Inches(2.0)
for title, desc in methods:
    add_text(slide, Inches(0.65), y, Inches(1.4), Inches(0.3),
             title, font_size=13, color=BLUE, bold=True)
    add_text(slide, Inches(2.1), y, Inches(4.3), Inches(0.6),
             desc, font_size=11, color=DARK, font_name="Consolas")
    y += Inches(0.7) if '\n' not in desc else Inches(0.85)

# Right top: Wire protocol
add_card(slide, Inches(6.9), Inches(1.45), Inches(6.0), Inches(2.7))
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.9), Inches(1.45), Inches(6.0), Inches(0.05))
bar.fill.solid(); bar.fill.fore_color.rgb = GREEN; bar.line.fill.background()
add_text(slide, Inches(7.1), Inches(1.55), Inches(5.6), Inches(0.4),
         "Wire Protocols (all BYTES payloads)", font_size=18, color=GREEN, bold=True)
add_code_block(slide, Inches(7.1), Inches(2.0), Inches(5.6), Inches(1.9),
    '// Text message\n'
    '"CHAT|senderName|timestamp|messageText"\n'
    '// GPS location\n'
    '"LOCATION|senderName|timestamp|lat|lng"\n'
    '// File (image, voice, etc.)\n'
    '"CHATFILE|senderName|fileName|mimeType|fileSize|"\n'
    '    + [raw file bytes]',
    font_size=12)

# Right bottom: State management
add_card(slide, Inches(6.9), Inches(4.4), Inches(6.0), Inches(2.65))
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.9), Inches(4.4), Inches(6.0), Inches(0.05))
bar.fill.solid(); bar.fill.fore_color.rgb = ORANGE; bar.line.fill.background()
add_text(slide, Inches(7.1), Inches(4.5), Inches(5.6), Inches(0.4),
         "State Management", font_size=18, color=ORANGE, bold=True)
add_bullet_list(slide, Inches(7.1), Inches(4.9), Inches(5.6), Inches(2.0), [
    "discoveredDevices / connectedDevices maps",
    "sendingFiles + openPfds (send-side tracking)",
    "pendingReceiveMeta + pendingReceivePayloads",
    "chatModeActive flag (invitation vs auto-accept)",
    "ExecutorService for background file I/O",
    "Handler(MainLooper) for UI-thread callbacks",
], font_size=12, color=DARK, bullet_color=ORANGE)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — TransferEventBus (Callback Interfaces)
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_LIGHT)
add_top_stripe(slide)

add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         "Event Bus & Callback Interfaces", font_size=36, color=BLACK, bold=True)
add_text(slide, Inches(0.6), Inches(0.85), Inches(12), Inches(0.4),
         "TransferEventBus.java  |  Decouples Network layer from UI layer  |  Observer pattern",
         font_size=15, color=GRAY)

interfaces = [
    ("ConnectionListener", BLUE,
     "interface ConnectionListener {\n"
     "    void onDevicesUpdated(List<DeviceInfo> devices);\n"
     "    void onConnectionEstablished(DeviceInfo device);\n"
     "    void onConnectionFailed(String reason);\n"
     "    void onDisconnected(String endpointId);\n"
     "    void onChatInvitationReceived(\n"
     "        String endpointId, String deviceName);\n"
     "}"),
    ("TransferListener", TEAL,
     "interface TransferListener {\n"
     "    void onIncomingFile(\n"
     "        FileMetadata meta, String fromDeviceName);\n"
     "    void onTransferProgressUpdated(\n"
     "        TransferProgress progress);\n"
     "    void onTransferCompleted(\n"
     "        TransferProgress finalProgress);\n"
     "    void onTransferFailed(\n"
     "        String fileName, String error);\n"
     "}"),
    ("ChatListener", ORANGE,
     "interface ChatListener {\n"
     "    void onChatMessageReceived(\n"
     "        String endpointId, String senderName,\n"
     "        String text, long timestamp);\n"
     "    void onLocationMessageReceived(\n"
     "        String endpointId, String senderName,\n"
     "        double lat, double lng,\n"
     "        long timestamp);\n"
     "}"),
]

for i, (name, color, code) in enumerate(interfaces):
    left = Inches(0.35 + i * 4.25)
    add_card(slide, left, Inches(1.45), Inches(4.05), Inches(3.3))
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(1.45), Inches(4.05), Inches(0.05))
    bar.fill.solid(); bar.fill.fore_color.rgb = color; bar.line.fill.background()
    add_text(slide, left + Inches(0.15), Inches(1.55), Inches(3.7), Inches(0.35),
             name, font_size=17, color=color, bold=True)
    add_code_block(slide, left + Inches(0.1), Inches(1.95), Inches(3.85), Inches(2.65),
                   code, font_size=11)

# Integration flow
add_card(slide, Inches(0.35), Inches(5.05), Inches(12.55), Inches(2.2))
add_text(slide, Inches(0.55), Inches(5.1), Inches(12), Inches(0.4),
         "Integration Flow", font_size=18, color=BLACK, bold=True)
add_code_block(slide, Inches(0.55), Inches(5.5), Inches(12.15), Inches(1.6),
    "// Device discovered\n"
    "NearbyConnectionsManager --> connectionListener.onDevicesUpdated() --> DeviceAdapter\n"
    "// User taps device\n"
    "nearbyMgr.connectToDevice() --> onConnectionEstablished() --> open ChatRoomFragment\n"
    "// User sends file (BYTES payload)\n"
    "nearbyMgr.sendFileInChat(BYTES) --> receiver.handleChatFileBytes() --> save to cache\n"
    "// File saved\n"
    "ChatFileReceivedCallback.onChatFileReceived(uri, meta) --> ChatAdapter.updateMessageUri()",
    font_size=11)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — RecyclerView Adapters
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_LIGHT)
add_top_stripe(slide)

add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         "RecyclerView Adapters (4)", font_size=36, color=BLACK, bold=True)
add_text(slide, Inches(0.6), Inches(0.85), Inches(12), Inches(0.4),
         "All extend RecyclerView.Adapter<ViewHolder>  |  ViewHolder pattern  |  Callback interfaces for clicks",
         font_size=15, color=GRAY)

adapters = [
    ("ChatAdapter", BLUE, "item_chat_message_sent.xml / received.xml", [
        "2 view types: VIEW_TYPE_SENT / RECEIVED",
        "Supports TEXT, FILE, LOCATION messages",
        "Inline ImageView for image thumbnails",
        "MediaPlayer for voice playback",
        "Dialog for full-screen image preview",
        "Open Map button for location",
        "updateMessageUri() for deferred file URIs",
        "SimpleDateFormat for timestamps",
    ]),
    ("DeviceAdapter", TEAL, "item_device.xml (MaterialCardView)", [
        "Device name + connection status",
        "OnDeviceClickListener callback",
        "Used in Discovery & Chat device list",
        "Chevron icon for visual affordance",
    ]),
    ("FileListAdapter", INDIGO, "item_shared_file.xml (MaterialCardView)", [
        "File name, size, uploader, timestamp",
        "OnFileActionListener (preview/download/delete)",
        "Used in FileListFragment (Shared Drive)",
        "Firebase real-time data source",
    ]),
    ("HistoryPeerAdapter", ORANGE, "item_history_peer.xml", [
        "Peer name + Reconnect button",
        "OnPeerClickListener callback",
        "OnReconnectClickListener callback",
        "Shows 'Connect' for online history peers",
    ]),
]

for i, (name, color, layout, items) in enumerate(adapters):
    col = i % 2
    row = i // 2
    left = Inches(0.4 + col * 6.35)
    top = Inches(1.35 + row * 3.05)
    add_card(slide, left, top, Inches(6.1), Inches(2.85))
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(6.1), Inches(0.05))
    bar.fill.solid(); bar.fill.fore_color.rgb = color; bar.line.fill.background()
    add_text(slide, left + Inches(0.15), top + Inches(0.08), Inches(3.5), Inches(0.35),
             name, font_size=18, color=color, bold=True)
    add_text(slide, left + Inches(0.15), top + Inches(0.38), Inches(5.8), Inches(0.3),
             "Layout: " + layout, font_size=11, color=GRAY, font_name="Consolas")
    for j, item in enumerate(items):
        add_text(slide, left + Inches(0.2), top + Inches(0.68 + j * 0.26), Inches(5.6), Inches(0.26),
                 "\u25B8  " + item, font_size=12, color=DARK)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Service & Storage
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_LIGHT)
add_top_stripe(slide)

add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         "Service & Storage", font_size=36, color=BLACK, bold=True)

# FileTransferService
add_card(slide, Inches(0.4), Inches(1.15), Inches(6.1), Inches(2.95))
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(1.15), Inches(6.1), Inches(0.05))
bar.fill.solid(); bar.fill.fore_color.rgb = BLUE; bar.line.fill.background()
add_text(slide, Inches(0.6), Inches(1.25), Inches(5.8), Inches(0.4),
         "FileTransferService", font_size=20, color=BLUE, bold=True)
add_code_block(slide, Inches(0.6), Inches(1.65), Inches(5.7), Inches(0.4),
    "public class FileTransferService extends LifecycleService", font_size=11)
add_bullet_list(slide, Inches(0.6), Inches(2.1), Inches(5.6), Inches(1.8), [
    "Foreground service (foregroundServiceType=dataSync)",
    "startForeground() within 5s (API 34+ requirement)",
    "NotificationCompat + NotificationChannel (API 26+)",
    "ExecutorService for background file I/O",
    "Parcelable extras (FileMetadata via Intent)",
], font_size=12, color=DARK, bullet_color=BLUE)

# FileStorageManager
add_card(slide, Inches(6.8), Inches(1.15), Inches(6.1), Inches(2.95))
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(1.15), Inches(6.1), Inches(0.05))
bar.fill.solid(); bar.fill.fore_color.rgb = TEAL; bar.line.fill.background()
add_text(slide, Inches(7.0), Inches(1.25), Inches(5.8), Inches(0.4),
         "FileStorageManager", font_size=20, color=TEAL, bold=True)
add_text(slide, Inches(7.0), Inches(1.6), Inches(5.8), Inches(0.3),
         "Scoped Storage adaptation by API level", font_size=13, color=GRAY)
add_bullet_list(slide, Inches(7.0), Inches(1.9), Inches(5.6), Inches(1.8), [
    "API 29+: MediaStore.Downloads + IS_PENDING",
    "API 24-28: Environment.DIRECTORY_DOWNLOADS",
    "ContentResolver.insert() + openOutputStream()",
    "Atomic publish (IS_PENDING 1 -> 0)",
    "Unique file naming with (1), (2)... suffixes",
], font_size=12, color=DARK, bullet_color=TEAL)

# ChatHistoryManager
add_card(slide, Inches(0.4), Inches(4.35), Inches(6.1), Inches(2.9))
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(4.35), Inches(6.1), Inches(0.05))
bar.fill.solid(); bar.fill.fore_color.rgb = ORANGE; bar.line.fill.background()
add_text(slide, Inches(0.6), Inches(4.45), Inches(5.8), Inches(0.4),
         "ChatHistoryManager", font_size=20, color=ORANGE, bold=True)
add_text(slide, Inches(0.6), Inches(4.8), Inches(5.8), Inches(0.3),
         "SharedPreferences + JSON serialization", font_size=13, color=GRAY)
add_bullet_list(slide, Inches(0.6), Inches(5.1), Inches(5.6), Inches(2.0), [
    'Key = "history_" + peerDeviceName',
    "JSONArray / JSONObject (manual serialization)",
    "Persists: type, sender, text, timestamp, fileMetadata, savedUri",
    "Max 200 messages per peer (FIFO trimming)",
    "updateFileUri() for deferred file URI updates",
], font_size=12, color=DARK, bullet_color=ORANGE)

# SharedDriveRepository
add_card(slide, Inches(6.8), Inches(4.35), Inches(6.1), Inches(2.9))
bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(4.35), Inches(6.1), Inches(0.05))
bar.fill.solid(); bar.fill.fore_color.rgb = INDIGO; bar.line.fill.background()
add_text(slide, Inches(7.0), Inches(4.45), Inches(5.8), Inches(0.4),
         "SharedDriveRepository", font_size=20, color=INDIGO, bold=True)
add_text(slide, Inches(7.0), Inches(4.8), Inches(5.8), Inches(0.3),
         "Firebase Cloud Storage + Firestore", font_size=13, color=GRAY)
add_bullet_list(slide, Inches(7.0), Inches(5.1), Inches(5.6), Inches(2.0), [
    "Two-phase upload: Storage -> Firestore metadata",
    "Real-time Firestore snapshot listener (no polling)",
    "UUID-based unique storage paths",
    "Upload progress callback",
    "Delete from both Storage and Firestore",
], font_size=12, color=DARK, bullet_color=INDIGO)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Model Classes
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_LIGHT)
add_top_stripe(slide)

add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         "Model Classes (5)", font_size=36, color=BLACK, bold=True)
add_text(slide, Inches(0.6), Inches(0.85), Inches(12), Inches(0.4),
         "Data layer  |  Parcelable for Intent passing  |  Wire serialization for Nearby Connections",
         font_size=15, color=GRAY)

models = [
    ("DeviceInfo", BLUE,
     "public class DeviceInfo {\n"
     "    String endpointId;\n"
     "    String deviceName;\n"
     "    boolean connected;\n"
     "}"),
    ("FileMetadata", TEAL,
     "public class FileMetadata\n"
     "    implements Parcelable {\n"
     "    String fileName, mimeType;\n"
     "    long fileSize;\n"
     '    // "FILE|name|mime|size"\n'
     "    byte[] toBytes();\n"
     "    static fromBytes(byte[]);\n"
     "}"),
    ("TransferProgress", ORANGE,
     "public class TransferProgress\n"
     "    implements Parcelable {\n"
     "    String fileName;\n"
     "    long bytesTransferred, totalBytes;\n"
     "    enum Status {\n"
     "        IN_PROGRESS, DONE, FAILED\n"
     "    }\n"
     "    int getProgressPercent();\n"
     "}"),
    ("ChatMessage", INDIGO,
     "public class ChatMessage {\n"
     "    enum Type { TEXT, FILE, LOCATION }\n"
     "    String senderName, text;\n"
     "    long timestamp;\n"
     "    boolean outgoing;\n"
     "    FileMetadata fileMetadata;\n"
     "    Uri savedUri;\n"
     "    double latitude, longitude;\n"
     "}"),
    ("SharedFile", PINK,
     "public class SharedFile {\n"
     '    @DocumentId String fileId;\n'
     "    String fileName, downloadUrl;\n"
     "    String mimeType, storagePath;\n"
     "    long fileSize;\n"
     '    @ServerTimestamp Date uploadedAt;\n'
     "    String getFileSizeFormatted();\n"
     "}"),
]

for i, (name, color, code) in enumerate(models):
    col = i % 3
    row = i // 3
    left = Inches(0.3 + col * 4.25)
    top = Inches(1.35 + row * 3.15)
    h = Inches(2.95) if row == 0 else Inches(2.85)
    add_card(slide, left, top, Inches(4.05), h)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, Inches(4.05), Inches(0.05))
    bar.fill.solid(); bar.fill.fore_color.rgb = color; bar.line.fill.background()
    add_text(slide, left + Inches(0.15), top + Inches(0.08), Inches(3.7), Inches(0.35),
             name, font_size=18, color=color, bold=True)
    add_code_block(slide, left + Inches(0.1), top + Inches(0.45), Inches(3.85), h - Inches(0.6),
                   code, font_size=11)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Permissions & API Level Adaptation
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_LIGHT)
add_top_stripe(slide)

add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         "Permissions & API Level Adaptation", font_size=36, color=BLACK, bold=True)
add_text(slide, Inches(0.6), Inches(0.85), Inches(12), Inches(0.4),
         "15 permissions declared  |  Runtime requests via ActivityResultLauncher  |  Graceful degradation",
         font_size=15, color=GRAY)

groups = [
    ("Nearby Connections", BLUE, [
        "BLUETOOTH, BLUETOOTH_ADMIN (API \u2264 30)",
        "BLUETOOTH_SCAN, BLUETOOTH_ADVERTISE,\n  BLUETOOTH_CONNECT (API 31+)",
        "NEARBY_WIFI_DEVICES (API 33+, neverForLocation)",
        "ACCESS_FINE_LOCATION (API 24-32)",
        "ACCESS_WIFI_STATE, CHANGE_WIFI_STATE",
    ]),
    ("Storage", TEAL, [
        "READ_EXTERNAL_STORAGE (API \u2264 32)",
        "WRITE_EXTERNAL_STORAGE (API \u2264 28)",
        "READ_MEDIA_IMAGES (API 33+)",
        "READ_MEDIA_VIDEO (API 33+)",
        "READ_MEDIA_FILES (API 33+)",
    ]),
    ("Other", ORANGE, [
        "INTERNET, ACCESS_NETWORK_STATE",
        "FOREGROUND_SERVICE (all APIs)",
        "FOREGROUND_SERVICE_DATA_SYNC (API 34+)",
        "POST_NOTIFICATIONS (API 33+)",
        "RECORD_AUDIO (voice messages)",
    ]),
]

for i, (title, color, perms) in enumerate(groups):
    left = Inches(0.3 + i * 4.25)
    add_card(slide, left, Inches(1.4), Inches(4.05), Inches(3.1))
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(1.4), Inches(4.05), Inches(0.05))
    bar.fill.solid(); bar.fill.fore_color.rgb = color; bar.line.fill.background()
    add_text(slide, left + Inches(0.15), Inches(1.5), Inches(3.7), Inches(0.35),
             title, font_size=17, color=color, bold=True)
    for j, perm in enumerate(perms):
        add_text(slide, left + Inches(0.2), Inches(1.9 + j * 0.4), Inches(3.6), Inches(0.4),
                 "\u25B8  " + perm, font_size=12, color=DARK)

# API adaptation code block
add_card(slide, Inches(0.3), Inches(4.8), Inches(12.7), Inches(2.4))
add_text(slide, Inches(0.5), Inches(4.85), Inches(12), Inches(0.35),
         "API Level Branching in Code", font_size=17, color=BLACK, bold=True)
add_code_block(slide, Inches(0.5), Inches(5.2), Inches(12.3), Inches(1.9),
    "// Storage\n"
    "if (Build.VERSION.SDK_INT >= 29)  MediaStore.Downloads + IS_PENDING  // scoped storage\n"
    "else                              Environment.DIRECTORY_DOWNLOADS     // direct write\n"
    "\n"
    "// Bluetooth permissions\n"
    "if (Build.VERSION.SDK_INT >= 31)  BLUETOOTH_SCAN, ADVERTISE, CONNECT  // granular\n"
    "else                              BLUETOOTH, BLUETOOTH_ADMIN           // legacy\n"
    "\n"
    "// Nearby discovery\n"
    "if (Build.VERSION.SDK_INT >= 33)  NEARBY_WIFI_DEVICES (neverForLocation)\n"
    "else                              ACCESS_FINE_LOCATION (required for Nearby)\n"
    "\n"
    "// API 36 workaround: Payload.fromFile() broken -> CHATFILE BYTES payload",
    font_size=12)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Components Summary Table
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_LIGHT)
add_top_stripe(slide)

add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         "Android Components Summary", font_size=36, color=BLACK, bold=True)

rows = [
    ("Component", "Count", "Techniques Used", True),
    ("Activity", "2", "AppCompatActivity, FragmentContainerView, BottomNavigationView, configChanges, Intent extras", False),
    ("Fragment", "5", "Fragment lifecycle, Bundle arguments, FragmentManager back stack, onDestroyView cleanup", False),
    ("Service", "1", "LifecycleService, Foreground service, NotificationChannel, Parcelable Intent extras", False),
    ("RecyclerView", "4", "ViewHolder pattern, 2 view types (ChatAdapter), MaterialCardView, click callbacks", False),
    ("ContentProvider", "1", "FileProvider for sharing cached files (API 24+)", False),
    ("Nearby Connections", "1", "P2P_CLUSTER, advertise/discover, BYTES payloads, PayloadCallback", False),
    ("Firebase", "1", "Firestore real-time listeners, Cloud Storage upload/download, UUID paths", False),
    ("SharedPreferences", "1", "JSON serialization, per-peer chat history, device name persistence", False),
    ("MediaRecorder", "-", "Audio recording (AAC/MPEG_4), temp file, sent via Nearby BYTES", False),
    ("MediaPlayer", "-", "Voice playback, prepareAsync + OnPreparedListener, error handling", False),
    ("LocationServices", "-", "FusedLocationProviderClient, GPS coordinates, shared via Nearby", False),
    ("OSMDroid", "1", "MapView, Marker/Polyline overlays, OSRM API routes, MyLocationOverlay", False),
]

y = Inches(1.15)
for i, (comp, count, tech, is_header) in enumerate(rows):
    h = Inches(0.38) if is_header else Inches(0.44)
    bg = BLUE if is_header else (BG_CARD if i % 2 == 1 else BG_LIGHT)
    txt_color = WHITE if is_header else BLACK
    tech_color = WHITE if is_header else DARK

    add_card(slide, Inches(0.4), y, Inches(12.5), h,
             fill=bg, border_color=BLUE if is_header else BORDER)
    add_text(slide, Inches(0.55), y + Inches(0.03), Inches(2.3), Inches(0.35),
             comp, font_size=13, color=txt_color, bold=True)
    add_text(slide, Inches(2.85), y + Inches(0.03), Inches(0.8), Inches(0.35),
             count, font_size=13, color=txt_color, bold=is_header, alignment=PP_ALIGN.CENTER)
    add_text(slide, Inches(3.7), y + Inches(0.03), Inches(9.0), Inches(0.35),
             tech, font_size=12 if not is_header else 13,
             color=tech_color, bold=is_header)
    y += h + Inches(0.02)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Class Hierarchy
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_LIGHT)
add_top_stripe(slide)

add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         "Class Hierarchy & Package Structure", font_size=36, color=BLACK, bold=True)

add_code_block(slide, Inches(0.4), Inches(1.05), Inches(12.5), Inches(6.1),
    "com.example.csci3310_airdrop_proj/\n"
    "|\n"
    "+-- MainActivity                       extends AppCompatActivity\n"
    "|      implements ConnectionListener, TransferListener, ChatListener,\n"
    "|                 OnFileSavedCallback, ChatFileReceivedCallback\n"
    "|\n"
    "+-- ui/\n"
    "|   +-- MapActivity                    extends AppCompatActivity\n"
    "|   +-- fragment/\n"
    "|   |   +-- FileListFragment           extends Fragment       // Shared Drive list\n"
    "|   |   +-- SendModeFragment           extends Fragment       // Upload to Firebase\n"
    "|   |   +-- DeviceDiscoveryFragment    extends Fragment       // Nearby device list\n"
    "|   |   +-- ChatRoomFragment           extends Fragment       // Chat UI + voice\n"
    "|   |   +-- ChatDeviceListFragment     extends Fragment       // Chat lobby\n"
    "|   +-- adapter/\n"
    "|       +-- ChatAdapter                extends RecyclerView.Adapter<ChatViewHolder>\n"
    "|       +-- DeviceAdapter              extends RecyclerView.Adapter<DeviceViewHolder>\n"
    "|       +-- FileListAdapter            extends RecyclerView.Adapter<FileViewHolder>\n"
    "|       +-- HistoryPeerAdapter         extends RecyclerView.Adapter<PeerViewHolder>\n"
    "|\n"
    "+-- network/\n"
    "|   +-- NearbyConnectionsManager       // Nearby Connections API wrapper\n"
    "|   +-- TransferEventBus               // 3 callback interfaces\n"
    "|\n"
    "+-- service/\n"
    "|   +-- FileTransferService            extends LifecycleService  // foreground\n"
    "|\n"
    "+-- storage/\n"
    "|   +-- FileStorageManager             // MediaStore / Downloads save\n"
    "|   +-- ChatHistoryManager             // SharedPreferences + JSON\n"
    "|\n"
    "+-- repository/\n"
    "|   +-- SharedDriveRepository          // Firebase Storage + Firestore\n"
    "|\n"
    "+-- model/\n"
    "    +-- DeviceInfo                     // POJO\n"
    "    +-- FileMetadata                   implements Parcelable  // toBytes/fromBytes\n"
    "    +-- TransferProgress               implements Parcelable  // status enum\n"
    "    +-- ChatMessage                    // POJO, Type enum\n"
    "    +-- SharedFile                     // @DocumentId, @ServerTimestamp",
    font_size=13)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Demo Flow
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_LIGHT)
add_top_stripe(slide)

add_text(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
         "Live Demo", font_size=36, color=BLACK, bold=True)

steps = [
    ("1", "Shared Drive", "Upload file to Firebase Storage \u2192 appears in real-time on all devices", INDIGO),
    ("2", "Device Discovery", "Both devices advertise + discover \u2192 device list populates automatically", BLUE),
    ("3", "Chat Connection", "Tap device \u2192 invitation dialog on receiver \u2192 Accept \u2192 chat room opens", TEAL),
    ("4", "Text Chat", "Send text messages \u2192 appear instantly on both sides via BYTES payloads", BLACK),
    ("5", "Image Sharing", "Send image \u2192 inline thumbnail preview \u2192 tap for fullscreen dialog", ORANGE),
    ("6", "Voice Message", "Hold record button \u2192 release \u2192 voice sent \u2192 Play button on receiver", PINK),
    ("7", "Location Sharing", "Tap location button \u2192 GPS coordinates sent \u2192 Open Map button", GREEN),
    ("8", "Map View", "OSMDroid map \u2192 current location overlay \u2192 walking route to shared location", TEAL),
    ("9", "Chat History", "Quit chat \u2192 reopen \u2192 all messages preserved (text + images + voice)", ORANGE),
]

for i, (num, title, desc, color) in enumerate(steps):
    y = Inches(1.05 + i * 0.68)
    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), y + Inches(0.04), Inches(0.4), Inches(0.4))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    add_text(slide, Inches(0.7), y + Inches(0.06), Inches(0.4), Inches(0.35),
             num, font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, Inches(1.3), y + Inches(0.04), Inches(2.3), Inches(0.4),
             title, font_size=17, color=color, bold=True)
    add_text(slide, Inches(3.7), y + Inches(0.07), Inches(9.2), Inches(0.35),
             desc, font_size=14, color=DARK)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Thank You
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, BG_WHITE)
# Accent rectangle at top
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(3.0))
shape.fill.solid()
shape.fill.fore_color.rgb = BLUE
shape.line.fill.background()

add_text(slide, Inches(1), Inches(0.8), Inches(11.3), Inches(1),
         "Thank You", font_size=60, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(1.9), Inches(11.3), Inches(0.6),
         "Questions?", font_size=32, color=RGBColor(0xBB, 0xDE, 0xFB), alignment=PP_ALIGN.CENTER)

add_text(slide, Inches(1), Inches(4.0), Inches(11.3), Inches(0.6),
         "DroidDrop  |  CSCI3310 Mobile Computing  |  Spring 2026",
         font_size=20, color=GRAY, alignment=PP_ALIGN.CENTER)


# ── Save ────────────────────────────────────────────────────────────────────
output_path = "/Users/yik/Documents/Uni/CSCI3310/CSCI3310_AirDrop_Proj/DroidDrop_Demo.pptx"
prs.save(output_path)
print(f"Saved to: {output_path}")
